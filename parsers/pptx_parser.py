from pptx import Presentation


def parse_pptx(file_path):
    """Parse a PPTX (PowerPoint Presentation) file and return a list of slides.

    Args:
        file_path (str): The path to the PPTX file to be parsed.

    Returns:
        List[Dict[str, str]]: A list of dictionaries, each containing the title,
        content of the slides, and any relevant metadata.
    """
    presentation = Presentation(file_path)
    slides = []

    for i, slide in enumerate(presentation.slides):
        content = ""
        for shape in slide.shapes:
            # Check if the shape has text and is not an empty placeholder
            if hasattr(shape, "text") and shape.text.strip():
                content += shape.text.strip() + "\n"

        slides.append({
            'title': f'Slide {i + 1}',
            'content': content.strip(),
            'slide_metadata': {}
        })

    return slides
