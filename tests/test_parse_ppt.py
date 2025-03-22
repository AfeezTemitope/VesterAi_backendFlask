import unittest
import os
from pptx import Presentation
from parsers.pptx_parser import parse_pptx


def create_sample_pptx(path):
    """Generate a PPTX file with sample content."""
    presentation = Presentation()

    slide1 = presentation.slides.add_slide(presentation.slide_layouts[5])  # Blank layout
    slide1.shapes.title.text = "Slide 1: Introduction"
    slide1.shapes.add_textbox(100, 100, 300, 100).text = "This is the first slide."

    slide2 = presentation.slides.add_slide(presentation.slide_layouts[5])
    slide2.shapes.title.text = "Slide 2: Testing PPTX"
    slide2.shapes.add_textbox(100, 100, 300, 100).text = "Second slide content here."
    presentation.save(path)


def create_invalid_pptx(path):
    """Create an invalid PPTX file."""
    with open(path, 'w') as f:
        f.write("This is not a valid PPTX file.")


class TestPPTXParser(unittest.TestCase):
    def setUp(self):
        """Set up test files."""
        self.sample_pptx_path = 'tests/sample.pptx'
        self.invalid_pptx_path = 'tests/invalid.pptx'

        create_sample_pptx(self.sample_pptx_path)
        create_invalid_pptx(self.invalid_pptx_path)

    def test_valid_pptx(self):
        """Test parsing a valid PPTX file."""
        slides = parse_pptx(self.sample_pptx_path)

        self.assertEqual(len(slides), 2)

        # Validate content of each slide
        self.assertEqual(slides[0]['title'], 'Slide 1')
        self.assertIn("Introduction", slides[0]['content'])
        self.assertEqual(slides[1]['title'], 'Slide 2')
        self.assertIn("Testing PPTX", slides[1]['content'])

    def test_invalid_pptx(self):
        """Test parsing an invalid PPTX file."""
        with self.assertRaises(Exception) as context:
            parse_pptx(self.invalid_pptx_path)
        self.assertIn("Package not found", str(context.exception))

    def tearDown(self):
        """Clean up the sample PPT files created."""
        os.remove(self.sample_pptx_path)
        os.remove(self.invalid_pptx_path)


if __name__ == '__main__':
    unittest.main()
